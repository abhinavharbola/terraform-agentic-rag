from pathlib import Path

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from bs4 import BeautifulSoup


def parse_pdf(path: Path) -> str:
    with pdfplumber.open(path) as pdf:
        return "\n\n".join(page.extract_text() or "" for page in pdf.pages)


def parse_docx(path: Path) -> str:
    doc = DocxDocument(path)
    return "\n".join(paragraph.text for paragraph in doc.paragraphs)


def parse_pptx(path: Path) -> str:
    presentation = Presentation(path)
    slides_text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                slides_text.append(shape.text_frame.text)
    return "\n".join(slides_text)


def parse_html(path: Path) -> str:
    soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
    return soup.get_text(separator="\n")


def parse_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".html": parse_html,
    ".htm": parse_html,
    ".txt": parse_txt,
    ".md": parse_txt,
}


def parse_file(path: Path) -> str:
    parser = PARSERS.get(path.suffix.lower())
    if parser is None:
        raise ValueError(f"no parser registered for extension: {path.suffix}")
    return parser(path)